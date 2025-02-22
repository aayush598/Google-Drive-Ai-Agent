from pydrive2.auth import GoogleAuth
from pydrive2.drive import GoogleDrive
import sqlite3
import io
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation



# Database path
db_path = "file_info.db"

def fetch_all_file_contents():
    """
    Fetches all readable files (TXT, PDF, DOCX, PPTX) from the database,
    downloads them from Google Drive, and extracts their content.
    Returns a list of dictionaries with file_name and content.
    """
    # Authenticate with Google Drive
    gauth = GoogleAuth()
    gauth.LocalWebserverAuth()  
    drive = GoogleDrive(gauth)

    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()

    # Query to fetch file details from the database
    cursor.execute("SELECT file_id, file_name, mime_type FROM files WHERE mime_type IN ('text/plain', 'application/pdf', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')")
    files = cursor.fetchall()

    file_contents = []

    for file_data in files:
        file_id, file_name, mime_type = file_data
        content = extract_content(file_id, file_name, mime_type,drive)
        file_contents.append({"file_name": file_name, "content": content})

    conn.close()
    return file_contents

def extract_content(file_id, file_name, mime_type,drive):
    """
    Downloads the file from Google Drive and extracts its content
    based on the file type.
    """
    try:
        # Download the file from Google Drive
        file = drive.CreateFile({'id': file_id})
        file.GetContentFile(file_name)

        # Extract content based on MIME type
        if mime_type == "text/plain":
            with open(file_name, "r", encoding="utf-8") as f:
                content = f.read()

        elif mime_type == "application/pdf":
            with open(file_name, "rb") as f:
                reader = PdfReader(f)
                content = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(file_name)
            content = "\n".join([para.text for para in doc.paragraphs])

        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(file_name)
            content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

        else:
            content = "Unsupported file type."

    except Exception as e:
        content = f"Error reading file: {str(e)}"

    return content
