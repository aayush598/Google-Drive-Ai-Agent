import sqlite3
from services.auth import drive
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# Database path
db_path = "file_info.db"

def fetch_all_file_contents():
    """
    Fetches all readable files (TXT, PDF, DOCX, PPTX) from the database,
    downloads them from Google Drive, and extracts their content.
    Returns a list of dictionaries with file_name and content.
    """
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    cursor.execute("""
        SELECT file_id, file_name, mime_type FROM files 
        WHERE mime_type IN (
            'text/plain', 'application/pdf', 
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
            'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    """)
    files = cursor.fetchall()
    conn.close()

    file_contents = []
    for file_id, file_name, mime_type in files:
        content = extract_content(file_id, file_name, mime_type, drive)
        file_contents.append({"file_name": file_name, "content": content})
    
    return file_contents

def extract_content(file_id, file_name, mime_type, drive):
    """
    Downloads the file from Google Drive and extracts its content
    based on the file type.
    """
    try:
        file = drive.CreateFile({'id': file_id})
        file.GetContentFile(file_name)

        if mime_type == "text/plain":
            with open(file_name, "r", encoding="utf-8") as f:
                return f.read()
        elif mime_type == "application/pdf":
            with open(file_name, "rb") as f:
                reader = PdfReader(f)
                return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(file_name)
            return "\n".join([para.text for para in doc.paragraphs])
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(file_name)
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        else:
            return "Unsupported file type."
    except Exception as e:
        return f"Error reading file: {str(e)}"
